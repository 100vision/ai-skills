import os
import sys
import json
import argparse
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt


class PPTSkillHandler:
    def __init__(self):
        self.current_template = "my_template.pptx"
        self.layout_idx_cache = {}

    def _load_registry(self) -> Optional[dict]:
        if os.path.exists("registry.json"):
            with open("registry.json", "r") as f:
                return json.load(f)
        return None

    def _resolve_template(self, template_key: str) -> Optional[str]:
        """Resolve template_key to an actual file path via registry.json."""
        registry = self._load_registry()
        if not registry:
            return None
        templates = registry.get("available_templates", {})
        info = templates.get(template_key)
        if info:
            return info["file_path"]
        # fallback to default
        default_key = registry.get("default_template")
        default_info = templates.get(default_key)
        if default_info:
            return default_info["file_path"]
        return None

    def inspect_and_build_engine(self, template_filename: str) -> dict:
        """功能 1.1 & 1.2：技能分析、侦察版式并自动建立映射关系"""
        if not os.path.exists(template_filename):
            return {"status": "error", "message": f"未找到模板文件: {template_filename}"}

        self.current_template = template_filename
        prs = Presentation(template_filename)

        for index, layout in enumerate(prs.slide_layouts):
            layout_name_lower = layout.name.lower()
            ph_map = {ph.name.lower(): ph.placeholder_format.idx for ph in layout.placeholders}

            # 智能特征识别匹配
            if "title" in layout_name_lower and "subtitle" in [p.name.lower() for p in layout.placeholders]:
                self.layout_idx_cache["cover"] = {
                    "layout_num": index,
                    "title_idx": self._find_key_idx(ph_map, "title"),
                    "sub_idx": self._find_key_idx(ph_map, "subtitle")
                }
            elif "section" in layout_name_lower or "header" in layout_name_lower:
                self.layout_idx_cache["transition"] = {"layout_num": index, "title_idx": next(iter(ph_map.values()), 0) if ph_map else 0}
            elif "three" in layout_name_lower or len(layout.placeholders) >= 4:
                sorted_ids = sorted(ph_map.values())
                self.layout_idx_cache["three_cols"] = {"layout_num": index, "title_idx": sorted_ids[0] if sorted_ids else 0, "col_ids": sorted_ids[1:4]}
            elif "two" in layout_name_lower or "comparison" in layout_name_lower:
                sorted_ids = sorted(ph_map.values())
                self.layout_idx_cache["text_only"] = {
                    "layout_num": index, "title_idx": sorted_ids[0] if sorted_ids else 0,
                    "left_idx": sorted_ids[1] if len(sorted_ids) > 1 else 0, "right_idx": sorted_ids[2] if len(sorted_ids) > 2 else 0
                }
                self.layout_idx_cache["img_text"] = {
                    "layout_num": index, "title_idx": sorted_ids[0] if sorted_ids else 0, "body_idx": sorted_ids[1] if len(sorted_ids) > 1 else 0
                }

        # 兜底安全机制：若模板命名不规范导致未识别，直接用物理序号 0,1,2,3,4 强行顶上
        defaults = ["cover", "transition", "text_only", "img_text", "three_cols"]
        for i, key in enumerate(defaults):
            if key not in self.layout_idx_cache and i < len(prs.slide_layouts):
                self.layout_idx_cache[key] = {"layout_num": i, "title_idx": 0, "sub_idx": 1, "left_idx": 1, "right_idx": 2, "body_idx": 1, "col_ids": [1, 2, 3]}

        return {"status": "success", "engine_cache": self.layout_idx_cache}

    def render_ppt_from_draft(self, payload_path: str, output_filename: str) -> dict:
        """功能 2：加载 AI 写入的 Payload 映射，驱动母版一键渲染成品"""
        if not os.path.exists(payload_path):
            return {"status": "error", "message": f"未找到 AI 写入的结构化载荷文件: {payload_path}"}

        with open(payload_path, "r") as f:
            payload_data = json.load(f)

        # ---- 多模板路由：从 registry 解析模板路径 ----
        template_key = payload_data.get("template_key")
        template_file = None
        cache_file = None

        if template_key:
            resolved = self._resolve_template(template_key)
            if resolved and os.path.exists(resolved):
                template_file = resolved
            cache_file = f".engine_cache_{template_key}.json"

        # 加载缓存：优先 per-template 缓存，fallback 到旧版 .engine_cache.json
        engine_cache = None
        if cache_file and os.path.exists(cache_file):
            with open(cache_file, "r") as f:
                engine_cache = json.load(f).get("engine_cache")
        elif os.path.exists(".engine_cache.json"):
            with open(".engine_cache.json", "r") as f:
                cache_data = json.load(f)
                engine_cache = cache_data.get("engine_cache")
                if not template_file:
                    template_file = cache_data.get("template_file", "my_template.pptx")

        if engine_cache:
            self.layout_idx_cache = engine_cache

        # 如果通过路由确定了模板但无缓存，自动静默分析
        if template_file and not engine_cache:
            if os.path.exists(template_file):
                res = self.inspect_and_build_engine(template_file)
                if res["status"] != "success":
                    return {"status": "error", "message": "自动模板分析失败"}
                # 保存 per-template 缓存
                if cache_file:
                    with open(cache_file, "w") as f:
                        json.dump({"engine_cache": self.layout_idx_cache}, f)
            else:
                return {"status": "error", "message": f"模板文件不存在: {template_file}"}

        # 最终兜底
        if not self.layout_idx_cache:
            return {"status": "error", "message": "错误：请先运行 inspect 模式扫描模板，生成引擎映射缓存。"}

        # ---- 渲染逻辑 ----
        slides_data = payload_data.get("slides", [])
        prs = Presentation(template_file or self.current_template)

        for slide in slides_data:
            l_type = slide.get("layout_type")
            cfg = self.layout_idx_cache.get(l_type, {"layout_num": 0, "title_idx": 0})

            current_slide = prs.slides.add_slide(prs.slide_layouts[cfg["layout_num"]])
            self._fill(current_slide, cfg.get("title_idx"), slide.get("title"))

            if l_type == "cover":
                self._fill(current_slide, cfg.get("sub_idx"), slide.get("subtitle", ""))
            elif l_type == "text_only":
                self._fill(current_slide, cfg.get("left_idx"), slide.get("left_text", ""))
                if slide.get("right_text"):
                    self._fill(current_slide, cfg.get("right_idx"), slide.get("right_text"))
            elif l_type == "img_text":
                self._fill(current_slide, cfg.get("body_idx"), slide.get("body_text", ""))
                img_p = slide.get("img_path", "placeholder.jpg")
                if os.path.exists(img_p):
                    current_slide.shapes.add_picture(img_p, Inches(1.0), Inches(2.0), width=Inches(5.0))
            elif l_type == "three_cols":
                col_ids = cfg.get("col_ids", [1, 2, 3])
                for idx_slot, col_key in enumerate(["col1", "col2", "col3"]):
                    col_data = slide.get(col_key, {"title": "", "desc": ""})
                    if idx_slot < len(col_ids):
                        ph = self._get_ph(current_slide, col_ids[idx_slot])
                        if ph:
                            tf = ph.text_frame
                            tf.text = col_data.get("title", "")
                            p = tf.add_paragraph()
                            p.text = col_data.get("desc", "")
                            p.font.size = Pt(14)

        prs.save(output_filename)
        return {"status": "success", "file_path": output_filename}

    def _find_key_idx(self, ph_map, keyword):
        for k, v in ph_map.items():
            if keyword in k:
                return v
        return next(iter(ph_map.values()), 0) if ph_map else 0

    def _get_ph(self, slide, target_idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == target_idx:
                return ph
        return None

    def _fill(self, slide, target_idx, text):
        ph = self._get_ph(slide, target_idx)
        if ph:
            ph.text = text


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument('--mode', choices=['inspect', 'render'], required=True)
    parser.add_argument('--file', type=str)
    parser.add_argument('--output', type=str, default='generated_presentation.pptx')
    args = parser.parse_args()

    handler = PPTSkillHandler()
    if args.mode == 'inspect':
        res = handler.inspect_and_build_engine(args.file)
        if res["status"] == "success":
            with open(".engine_cache.json", "w") as f:
                json.dump({"template_file": args.file, "engine_cache": res["engine_cache"]}, f)
        print(json.dumps(res))
    elif args.mode == 'render':
        res = handler.render_ppt_from_draft(".ppt_payload.json", args.output)
        print(json.dumps(res))
